"""
Document ingestion pipeline for the Academic Knowledge Intelligence Layer.

Pipeline:
  Upload → Format Detection → Text Extraction → Cleaning →
  Semantic Chunking → Academic Classification → Storage

Supported formats: PDF, DOCX, PPTX, TXT, MD, PNG, JPG, JPEG
"""

from __future__ import annotations

import io
import logging
import re
import time
import base64
from pathlib import Path
from typing import Any
from uuid import uuid4

from sqlalchemy.orm import Session

from ..config import settings
from .chunking import AcademicChunk, semantic_chunk, count_tokens, infer_microchunk_type
from .classifier import classify_chunk
from ..llm_pipeline import LLMCall
from ..models import Question

from .models import (
    AcademicDocument,
    ChunkApprovalStatus,
    DocumentType,
    KnowledgeChunk,
    ProcessingStatus,
    SubjectSyllabus,
)

logger = logging.getLogger("app.academic.ingestion")


# ---------------------------------------------------------------------------
# Text Extraction
# ---------------------------------------------------------------------------

def extract_text_from_pdf(content: bytes) -> tuple[str, int, dict[int, int]]:
    """Extract text from PDF with page tracking.
    
    Returns: (full_text, page_count, char_offset_to_page_map)
    """
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages_text: list[str] = []
    page_offsets: dict[int, int] = {}
    offset = 0

    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        page_offsets[offset] = page_num
        pages_text.append(text)
        offset += len(text) + 1  # +1 for newline
        
        # Optional: Extract images if PyMuPDF is available
        # This writes to image_chunks.json asynchronously or silently handles errors
        try:
            _extract_and_store_images_pymupdf(content, page_num)
        except Exception as e:
            logger.debug(f"Image extraction skipped or failed: {e}")

    full_text = "\n".join(pages_text)
    return full_text, len(reader.pages), page_offsets

def _extract_and_store_images_pymupdf(content: bytes, page_num: int):
    """
    Helper to extract images using PyMuPDF (fitz) and store in image_chunks.json.
    Very lightweight and isolated.
    """
    import fitz # PyMuPDF
    import json
    import os
    
    doc = fitz.open(stream=content, filetype="pdf")
    page = doc.load_page(page_num - 1)
    images = page.get_images(full=True)
    
    if not images:
        return
        
    chunks_path = Path("image_chunks.json")
    chunks = []
    if chunks_path.exists():
        try:
            with open(chunks_path, "r") as f:
                chunks = json.load(f)
        except json.JSONDecodeError:
            chunks = []
            
    images_dir = Path("static/images")
    images_dir.mkdir(parents=True, exist_ok=True)
    
    for img_index, img in enumerate(images):
        xref = img[0]
        base_image = doc.extract_image(xref)
        image_bytes = base_image["image"]
        image_ext = base_image["ext"]
        
        # Avoid tiny icons
        if len(image_bytes) < 10000:
            continue
            
        image_filename = f"img_{uuid4().hex[:8]}.{image_ext}"
        image_path = images_dir / image_filename
        
        with open(image_path, "wb") as f:
            f.write(image_bytes)
            
        # Extract surrounding text as proxy for VLM
        text_blocks = page.get_text("blocks")
        nearby_text = " ".join([b[4] for b in text_blocks[:3]]).strip().replace("\n", " ")
        
        chunks.append({
            "image_path": str(image_path),
            "page_num": page_num,
            "topic": nearby_text[:100] if nearby_text else "Academic Diagram",
            "caption": "Extracted academic diagram",
            "keywords": ["diagram", "academic", "figure"]
        })
        
    with open(chunks_path, "w") as f:
        json.dump(chunks, f, indent=2)


def extract_text_from_docx(content: bytes) -> tuple[str, int]:
    """Extract text from DOCX files including tables.
    
    Returns: (full_text, estimated_page_count)
    """
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(content))
    parts: list[str] = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)

    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                parts.append(" | ".join(row_texts))

    full_text = "\n".join(parts)
    # Rough page estimate: ~500 words per page
    est_pages = max(1, len(full_text.split()) // 500)
    return full_text, est_pages


def extract_text_from_pptx(content: bytes) -> tuple[str, int]:
    """Extract text from PPTX files.
    
    Returns: (full_text, slide_count)
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed, cannot extract PPTX text")
        return "", 0

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text += text + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        slide_text += " | ".join(row_texts) + "\n"
        parts.append(slide_text)

    return "\n".join(parts), len(prs.slides)


def extract_text_from_image(content: bytes) -> str:
    """Extract text from images using OCR.
    
    Tries local OCR first, then falls back to the configured vision model.
    """
    try:
        import pytesseract
        from PIL import Image
        image = Image.open(io.BytesIO(content))
        text = pytesseract.image_to_string(image).strip()
        if text:
            return text
    except ImportError:
        logger.info("pytesseract not installed, image OCR unavailable")
    except Exception as e:
        logger.warning("Image OCR failed: %s", e)

    llm = LLMCall(
        model=settings.ollama_vision_model,
        timeout=settings.ollama_request_timeout_seconds,
    )
    if not llm.is_available():
        return ""

    system = (
        "You are performing OCR for academic material ingestion. "
        "Transcribe the visible text faithfully and do not summarize."
    )
    prompt = (
        "Read the academic image and return only the extracted text. "
        "Preserve headings, lists, and technical terms whenever possible."
    )
    encoded = base64.b64encode(content).decode("utf-8")
    extracted = llm.generate_text(prompt, system, images=[encoded], model=settings.ollama_vision_model)
    return extracted.strip() if extracted else ""


def extract_text(file_name: str, content: bytes) -> tuple[str, int, dict[int, int] | None]:
    """
    Detect file format and extract text.
    
    Returns: (text, page_count, page_offset_map_or_None)
    """
    suffix = Path(file_name).suffix.lower()
    
    if suffix == ".pdf":
        text, pages, offsets = extract_text_from_pdf(content)
        return text, pages, offsets
    elif suffix == ".docx":
        text, pages = extract_text_from_docx(content)
        return text, pages, None
    elif suffix == ".pptx":
        text, slides = extract_text_from_pptx(content)
        return text, slides, None
    elif suffix in {".txt", ".md", ".csv"}:
        text = content.decode("utf-8", errors="ignore")
        pages = max(1, len(text.split()) // 500)
        return text, pages, None
    elif suffix in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}:
        text = extract_text_from_image(content)
        return text, 1, None
    else:
        text = content.decode("utf-8", errors="ignore")
        return text, 1, None


# ---------------------------------------------------------------------------
# Text Cleaning
# ---------------------------------------------------------------------------

def clean_academic_text(text: str) -> str:
    """Clean extracted text while preserving academic structure."""
    if not text:
        return ""

    # Normalize whitespace within lines (but preserve paragraph breaks)
    lines = text.split("\n")
    cleaned_lines: list[str] = []
    
    for line in lines:
        # Collapse multiple spaces
        line = re.sub(r"[ \t]+", " ", line)
        # Remove control characters except newline
        line = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", line)
        cleaned_lines.append(line.strip())

    result = "\n".join(cleaned_lines)
    # Collapse more than 3 consecutive newlines
    result = re.sub(r"\n{4,}", "\n\n\n", result)
    return result.strip()


# ---------------------------------------------------------------------------
# Document Type Detection
# ---------------------------------------------------------------------------

def detect_document_type(file_name: str, text: str) -> DocumentType:
    """Heuristically detect the type of academic document."""
    name_lower = file_name.lower()
    text_lower = text[:3000].lower() if text else ""

    if any(kw in name_lower for kw in ("syllabus", "curriculum")):
        return DocumentType.SYLLABUS
    if any(kw in name_lower for kw in ("previous", "model paper", "past paper", "question paper")):
        return DocumentType.PREVIOUS_PAPER
    if any(kw in name_lower for kw in ("question bank", "qbank", "q-bank")):
        return DocumentType.QUESTION_BANK
    if any(kw in name_lower for kw in ("lab", "manual", "experiment")):
        return DocumentType.LAB_MANUAL
    if name_lower.endswith(".pptx") or name_lower.endswith(".ppt"):
        return DocumentType.PPT

    # Content-based detection
    if any(kw in text_lower for kw in ("syllabus", "course objectives", "course outcomes")):
        return DocumentType.SYLLABUS
    if re.search(r"(?:Q\.?\s*\d|question\s+paper|marks?:?\s*\d)", text_lower):
        return DocumentType.QUESTION_BANK

    return DocumentType.NOTES


# ---------------------------------------------------------------------------
# Main Ingestion Pipeline
# ---------------------------------------------------------------------------

def create_document_record(
    db: Session,
    subject_id: int,
    user_id: int,
    file_name: str,
    content: bytes,
    document_type: DocumentType | None = None,
) -> AcademicDocument:
    """Save the file and create the initial database record immediately."""
    suffix = Path(file_name).suffix.lower()

    # --- Save file ---
    upload_dir = settings.storage_path / "academic" / str(subject_id)
    upload_dir.mkdir(parents=True, exist_ok=True)
    saved_path = upload_dir / f"{uuid4().hex}_{file_name}"
    saved_path.write_bytes(content)

    # --- Create document record ---
    doc = AcademicDocument(
        subject_id=subject_id,
        uploaded_by=user_id,
        file_name=file_name,
        file_type=suffix.lstrip("."),
        document_type=document_type or DocumentType.NOTES,
        storage_path=str(saved_path),
        processing_status=ProcessingStatus.EXTRACTING,
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)
    return doc

def pre_generate_topic_variants(
    db: Session,
    doc_id: int,
    subject_id: int,
    teacher_id: int,
    db_chunks: list[KnowledgeChunk],
) -> list[Question]:
    """Pre-generate a rich set of question variants (5M, 10M, split, image, definition, etc.)
    for every unique topic discovered in the ingested document.
    This guarantees high coverage and stable allocation during the review stage.
    """
    from .templates import compile_question
    from ..models import Question

    # Group by unique topic_name
    topics_map = {}
    for chunk in db_chunks:
        if not chunk.topic_name or chunk.topic_name.lower() == "unclassified":
            continue
        topic = chunk.topic_name
        if topic not in topics_map:
            topics_map[topic] = {
                "module": chunk.module_number or 1,
                "co": chunk.co_mapping or "CO1",
                "bloom": chunk.bloom_level or "L2",
                "keywords": [chunk.topic_name],
            }

    generated_questions = []
    used_templates = set()

    # Define the variants to generate per topic
    # Format: (suffix_name, bloom_level, marks, is_image, tag)
    variant_specs = [
        # 5M variants (3-5 variants)
        ("5M_v1", "L2", 5, False, "variant:5M"),
        ("5M_v2", "L2", 5, False, "variant:5M"),
        ("5M_v3", "L3", 5, False, "variant:5M"),
        # 10M variants (3-5 variants)
        ("10M_v1", "L2", 10, False, "variant:10M"),
        ("10M_v2", "L3", 10, False, "variant:10M"),
        ("10M_v3", "L4", 10, False, "variant:10M"),
        # split variants (2-3 variants)
        ("split_v1", "L2", 5, False, "variant:split"),
        ("split_v2", "L3", 5, False, "variant:split"),
        # application (2 variants)
        ("app_v1", "L3", 5, False, "variant:application"),
        ("app_v2", "L3", 10, False, "variant:application"),
        # analytical (2 variants)
        ("analysis_v1", "L4", 5, False, "variant:analysis"),
        ("analysis_v2", "L4", 10, False, "variant:analysis"),
        # image-based (1-2 variants)
        ("image_v1", "L2", 5, True, "variant:image"),
        ("image_v2", "L3", 10, True, "variant:image"),
        # definition / short-answer
        ("def_v1", "L1", 2, False, "variant:definition"),
        ("def_v2", "L1", 5, False, "variant:definition"),
    ]

    for topic_name, meta in topics_map.items():
        module = meta["module"]
        co = meta["co"]
        keywords = meta["keywords"]

        for suffix, bloom, marks, is_image, variant_tag in variant_specs:
            try:
                text = compile_question(
                    topic=topic_name,
                    bloom_level=bloom,
                    keywords=keywords,
                    marks=marks,
                    is_image_question=is_image,
                    used_templates=used_templates,
                )
                q = Question(
                    subject_id=subject_id,
                    teacher_id=teacher_id,
                    source_doc_id=doc_id,
                    text=text,
                    marks=marks,
                    course_outcome=co,
                    bloom_level=bloom,
                    difficulty="balanced" if marks <= 5 else "hard",
                    module_number=module,
                    tags=["pre_generated", "variant_compiler", variant_tag, f"topic:{topic_name}"],
                    is_verified=True,  # pre-generated variants are verified by template compile
                )
                generated_questions.append(q)
            except Exception as e:
                logger.warning(
                    "Failed compiling variant %s for topic %s: %s",
                    suffix,
                    topic_name,
                    e,
                )

    return generated_questions


def process_document_background(
    document_id: int,
    auto_approve_threshold: float = 0.6,
) -> None:
    """
    Background worker for extraction, chunking, classification, and embedding.
    Runs completely separate from the request thread.
    """
    from ..database import SessionLocal
    db = SessionLocal()
    
    try:
        doc = db.get(AcademicDocument, document_id)
        if not doc:
            return

        start_time = time.time()
        file_name = doc.file_name

        try:
            content = Path(doc.storage_path).read_bytes()

            # --- Extract text ---
            text, page_count, page_offsets = extract_text(file_name, content)
            doc.page_count = page_count

            if not text or len(text.strip()) < 50:
                doc.processing_status = ProcessingStatus.FAILED
                doc.processing_error = "Insufficient text extracted from document"
                db.commit()
                return

            # --- Clean text ---
            cleaned = clean_academic_text(text)
            doc.extracted_text = cleaned[:100000]  # Store first 100k chars

            # --- Multimodal structured parsing ---
            doc.processing_status = ProcessingStatus.PARSING
            db.commit()

            try:
                from .multimodal import parse_document_structure
                from ..config import settings

                struct_doc = parse_document_structure(
                    file_name=file_name,
                    content=content,
                    storage_dir=settings.storage_path,
                    analyze_figures=True,
                    extract_math=True,
                )

                # Persist structured content
                doc.structured_content = struct_doc.to_dict()
                doc.has_equations = struct_doc.has_equations
                doc.has_figures = struct_doc.has_figures
                doc.has_tables = struct_doc.has_tables
                doc.equation_count = struct_doc.summary.get("equation_count", 0)
                doc.figure_count = struct_doc.summary.get("figure_count", 0)
                doc.table_count = struct_doc.summary.get("table_count", 0)
                db.commit()

                # Replace plain cleaned text with enriched text for chunking
                # The enriched text embeds [EQUATION:...] and [FIGURE:...] markers
                # so the RAG system can reference visual content in question generation
                enriched = struct_doc.to_enriched_text()
                if enriched and len(enriched.strip()) > 50:
                    cleaned = enriched
                    logger.info(
                        "Multimodal parsing enriched text for '%s': "
                        "%d equations, %d figures, %d tables",
                        file_name,
                        doc.equation_count,
                        doc.figure_count,
                        doc.table_count,
                    )
            except Exception as exc:
                logger.warning(
                    "Multimodal structured parsing failed for '%s', using plain text: %s",
                    file_name, exc,
                )
                # Not fatal — continue with plain text

            # --- Detect document type ---
            if doc.document_type == DocumentType.NOTES:
                doc.document_type = detect_document_type(file_name, cleaned)

            # --- Semantic chunking ---
            doc.processing_status = ProcessingStatus.CHUNKING
            db.commit()

            chunks = semantic_chunk(
                cleaned,
                min_tokens=100,
                max_tokens=250,
                overlap_ratio=0.12,
                page_numbers=page_offsets,
            )

            if not chunks:
                logger.warning(
                    "Semantic chunking produced no chunks for '%s'; creating a fallback chunk",
                    file_name,
                )
                fallback_text = cleaned[:4000].strip()
                if not fallback_text:
                    doc.processing_status = ProcessingStatus.FAILED
                    doc.processing_error = "No meaningful chunks could be created"
                    db.commit()
                    return
                chunks = [
                    AcademicChunk(
                        text=fallback_text,
                        chunk_index=0,
                        token_count=count_tokens(fallback_text),
                        page_number=1,
                        source_section=None,
                    )
                ]

            # --- Load syllabus for classification ---
            syllabus = db.query(SubjectSyllabus).filter_by(subject_id=doc.subject_id).first()
            syllabus_modules = syllabus.modules_json if syllabus else None

            # --- Classify, generate questions, and store chunks concurrently ---
            doc.processing_status = ProcessingStatus.EMBEDDING
            db.commit()

            import asyncio
            import base64
            from ..llm_pipeline import get_pipeline
            from ..models import Question

            async def run_ingestion_parallel(chunks, syllabus_modules, struct_doc_opt):
                loop = asyncio.get_running_loop()
                pipeline = get_pipeline()
                
                async def process_chunk(chunk):
                    class_task = loop.run_in_executor(
                        None, classify_chunk, chunk.text, chunk.source_section, syllabus_modules
                    )
                    gen_task = loop.run_in_executor(
                        None, pipeline.extractor.from_text, chunk.text
                    )
                    try:
                        classification, extracted_qs = await asyncio.gather(class_task, gen_task)
                    except Exception as e:
                        logger.error(f"Chunk processing failed: {e}")
                        return chunk, None, []
                    return chunk, classification, extracted_qs
                
                chunk_tasks = [process_chunk(c) for c in chunks]
                
                async def process_image(fb):
                    if not fb.image_path:
                        return fb, []
                    def encode_and_extract():
                        try:
                            with open(fb.image_path, "rb") as img_file:
                                b64 = base64.b64encode(img_file.read()).decode("utf-8")
                            return pipeline.extractor.from_image(b64)
                        except Exception as e:
                            logger.error(f"Image extraction failed for {fb.image_path}: {e}")
                            return []
                    extracted_qs = await loop.run_in_executor(None, encode_and_extract)
                    return fb, extracted_qs
                
                image_tasks = []
                if struct_doc_opt:
                    image_tasks = [process_image(fb) for fb in struct_doc_opt.get_all_blocks("figure")]
                    
                chunk_results, image_results = await asyncio.gather(
                    asyncio.gather(*chunk_tasks),
                    asyncio.gather(*image_tasks)
                )
                
                return chunk_results, image_results
                
            struct_doc_opt = struct_doc if 'struct_doc' in locals() else None
            
            import threading
            
            def run_async(coro):
                try:
                    loop = asyncio.get_running_loop()
                except RuntimeError:
                    return asyncio.run(coro)
                
                result = [None]
                exception = [None]
                
                def target():
                    try:
                        new_loop = asyncio.new_event_loop()
                        asyncio.set_event_loop(new_loop)
                        result[0] = new_loop.run_until_complete(coro)
                        new_loop.close()
                    except Exception as e:
                        exception[0] = e
                        
                thread = threading.Thread(target=target)
                thread.start()
                thread.join()
                if exception[0]:
                    raise exception[0]
                return result[0]
                
            chunk_results, image_results = run_async(run_ingestion_parallel(chunks, syllabus_modules, struct_doc_opt))

            db_chunks = []
            db_questions = []
            last_detected_module = None
            
            for chunk, classification, extracted_qs in chunk_results:
                if not classification:
                    continue
                module_number = classification.module_number
                
                if module_number is None and syllabus_modules and chunk.page_number:
                    for mod in syllabus_modules:
                        if isinstance(mod, dict):
                            ps = mod.get("page_start")
                            pe = mod.get("page_end")
                            if ps and pe and ps <= chunk.page_number <= pe:
                                module_number = mod.get("module")
                                break

                if module_number is not None:
                    last_detected_module = module_number
                else:
                    module_number = last_detected_module

                approval = ChunkApprovalStatus.AUTO_APPROVED

                db_chunk = KnowledgeChunk(
                    document_id=doc.id,
                    subject_id=doc.subject_id,
                    chunk_text=chunk.text,
                    chunk_summary=infer_microchunk_type(chunk.text),
                    chunk_index=chunk.chunk_index,
                    token_count=chunk.token_count,
                    module_number=module_number,
                    syllabus_unit=None,
                    topic_name=classification.topic_name,
                    bloom_level=classification.bloom_level,
                    co_mapping=classification.co_mapping,
                    page_number=chunk.page_number,
                    confidence_score=classification.confidence_score,
                    approval_status=approval,
                )
                db.add(db_chunk)
                db_chunks.append(db_chunk)
                
                for q in extracted_qs:
                    q_mod = q.get("module") or module_number or 1
                    try:
                        marks = int(q.get("marks", 5))
                    except (ValueError, TypeError):
                        marks = 5
                    db_questions.append(Question(
                        subject_id=doc.subject_id,
                        teacher_id=doc.uploaded_by,
                        source_doc_id=doc.id,
                        text=q.get("text", "Unknown"),
                        marks=marks,
                        course_outcome=classification.co_mapping or "CO1",
                        bloom_level=q.get("bloom_level", classification.bloom_level or "L2"),
                        difficulty=q.get("difficulty", "medium"),
                        module_number=q_mod,
                        tags=["pre_generated", "text_extraction"],
                        is_verified=False
                    ))

            doc.total_chunks = len(chunks)
            db.commit()

            # Image-to-Vector Linkage
            for fb, extracted_qs in image_results:
                fig_desc = fb.analysis.get("description", "") if fb.analysis else ""
                fig_text = f"[FIGURE_PATH: {fb.image_path}]\\nCaption: {fb.caption}\\nDescription: {fig_desc}"
                fig_chunk = KnowledgeChunk(
                    document_id=doc.id,
                    subject_id=doc.subject_id,
                    chunk_text=fig_text,
                    chunk_index=len(db_chunks),
                    token_count=count_tokens(fig_text),
                    module_number=last_detected_module,
                    topic_name=(fb.caption[:250] if fb.caption else "Figure"),
                    bloom_level="L2",
                    co_mapping="CO2",
                    page_number=fb.page_number,
                    confidence_score=0.9,
                    approval_status=ChunkApprovalStatus.AUTO_APPROVED,
                )
                db.add(fig_chunk)
                db_chunks.append(fig_chunk)
                
                for q in extracted_qs:
                    try:
                        marks = int(q.get("marks", 5))
                    except (ValueError, TypeError):
                        marks = 5
                    db_questions.append(Question(
                        subject_id=doc.subject_id,
                        teacher_id=doc.uploaded_by,
                        source_doc_id=doc.id,
                        text=q.get("text", "Unknown"),
                        marks=marks,
                        course_outcome="CO2",
                        bloom_level=q.get("bloom_level", "L3"),
                        difficulty=q.get("difficulty", "medium"),
                        module_number=last_detected_module or 1,
                        tags=["pre_generated", "image_extraction"],
                        image_path=fb.image_path,
                        is_verified=False
                    ))
            
            db.add_all(db_questions)
            db.commit()

            # --- PRE-GENERATE MULTIPLE VARIANTS PER TOPIC ---
            try:
                variants = pre_generate_topic_variants(
                    db=db,
                    doc_id=doc.id,
                    subject_id=doc.subject_id,
                    teacher_id=doc.uploaded_by,
                    db_chunks=db_chunks,
                )
                if variants:
                    db.add_all(variants)
                    db.commit()
                    # Append to db_questions so their embeddings are also generated!
                    db_questions.extend(variants)
                    logger.info(
                        "Pre-generated %d topic variants for document %d",
                        len(variants),
                        doc.id,
                    )
            except Exception as ev:
                logger.error("Failed pre-generating topic variants: %s", ev)

            elapsed = time.time() - start_time
            logger.info(
                "Ingested '%s': %d pages, %d chunks, %d questions in %.2fs",
                file_name, page_count, len(db_chunks), len(db_questions), elapsed,
            )

            # --- Generate Embeddings ---
            from .embeddings import generate_embeddings_batch
            
            texts = [c.chunk_text for c in db_chunks]
            embeddings = generate_embeddings_batch(texts)
            for chunk, embedding in zip(db_chunks, embeddings):
                if embedding is not None:
                    chunk.embedding_vector = embedding
                    
            q_texts = [q.text for q in db_questions]
            if q_texts:
                q_embeddings = generate_embeddings_batch(q_texts)
                for q, embedding in zip(db_questions, q_embeddings):
                    if embedding is not None:
                        q.embedding_vector = embedding
                        
            doc.extracted_text = None  # Discard full extracted text as per guidelines

            doc.processing_status = ProcessingStatus.COMPLETED
            db.commit()
            
            # Dynamically regenerate Course Outcome (CO) descriptions based on newly ingested content!
            try:
                from .co_description_generator import generate_subject_co_descriptions
                generate_subject_co_descriptions(db, doc.subject_id)
                logger.info("Automatically generated dynamic CO descriptions for subject %d", doc.subject_id)
            except Exception as eco:
                logger.error("Failed to automatically generate dynamic CO descriptions: %s", eco)

            logger.info("Completed embedding generation for doc %d synchronously in background thread", doc.id)

        except Exception as e:
            logger.error("Ingestion failed for '%s': %s", file_name, e)
            doc.processing_status = ProcessingStatus.FAILED
            doc.processing_error = str(e)[:2000]
            db.commit()

    finally:
        db.close()
