"""
Structured Document Parser — Main Orchestrator.

Converts academic documents into a structured JSON representation:

{
  "pages": [
    {
      "page_number": 1,
      "blocks": [
        {"type": "heading",   "content": "Transformer Architecture"},
        {"type": "paragraph", "content": "Transformers use..."},
        {"type": "equation",  "latex": "\\alpha = softmax(QK^T)", "raw": "a = softmax qkt"},
        {"type": "figure",    "caption": "Fig 1: Encoder block",
                              "image_path": "...",
                              "analysis": {...}},
        {"type": "table",     "content": "...", "rows": [...]},
      ]
    }
  ],
  "summary": {
    "total_pages": 5,
    "has_equations": true,
    "has_figures": true,
    "has_tables": true,
    "equation_count": 3,
    "figure_count": 2,
    "table_count": 1,
    "heading_count": 8,
  }
}

Entry point: parse_document_structure(file_name, content, storage_dir)
"""

from __future__ import annotations

import io
import json
import logging
import os
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Any
from uuid import uuid4

logger = logging.getLogger("app.academic.multimodal.structured_parser")


# ---------------------------------------------------------------------------
# Data models
# ---------------------------------------------------------------------------

@dataclass
class PageBlock:
    """A single content block on a page."""
    type: str               # heading | paragraph | equation | figure | table | list | code | caption
    content: str = ""       # primary text content
    # Equation-specific
    latex: str = ""
    latex_method: str = ""  # nougat | vision_llm | heuristic
    # Figure-specific
    caption: str = ""
    image_path: str = ""    # relative path to saved image file
    analysis: dict[str, Any] = field(default_factory=dict)
    # Table-specific
    rows: list[list[str]] = field(default_factory=list)
    # Common metadata
    confidence: float = 1.0
    page_number: int = 0
    bbox: list[float] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        d: dict[str, Any] = {"type": self.type}
        if self.content:
            d["content"] = self.content
        if self.latex:
            d["latex"] = self.latex
            d["latex_method"] = self.latex_method
        if self.caption:
            d["caption"] = self.caption
        if self.image_path:
            d["image_path"] = self.image_path
        if self.analysis:
            d["analysis"] = self.analysis
        if self.rows:
            d["rows"] = self.rows
        if self.bbox:
            d["bbox"] = self.bbox
        d["confidence"] = round(self.confidence, 3)
        return d


@dataclass
class StructuredPage:
    page_number: int
    blocks: list[PageBlock] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "page_number": self.page_number,
            "blocks": [b.to_dict() for b in self.blocks],
        }


@dataclass
class StructuredDocument:
    """Complete structured representation of a parsed document."""
    pages: list[StructuredPage] = field(default_factory=list)
    summary: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        return {
            "pages": [p.to_dict() for p in self.pages],
            "summary": self.summary,
        }

    def to_json(self) -> str:
        return json.dumps(self.to_dict(), ensure_ascii=False, indent=2)

    @property
    def has_equations(self) -> bool:
        return self.summary.get("has_equations", False)

    @property
    def has_figures(self) -> bool:
        return self.summary.get("has_figures", False)

    @property
    def has_tables(self) -> bool:
        return self.summary.get("has_tables", False)

    def get_all_blocks(self, block_type: str | None = None) -> list[PageBlock]:
        """Get all blocks across all pages, optionally filtered by type."""
        blocks = [b for page in self.pages for b in page.blocks]
        if block_type:
            blocks = [b for b in blocks if b.type == block_type]
        return blocks

    def to_enriched_text(self) -> str:
        """
        Convert structured document to enriched text suitable for RAG chunking.

        Equations are represented as [EQ: latex], figures as [FIG: description],
        tables as [TABLE: content]. This preserves semantic information in text form.
        """
        parts: list[str] = []
        for page in self.pages:
            parts.append(f"\n--- Page {page.page_number} ---")
            for block in page.blocks:
                if block.type == "heading":
                    parts.append(f"\n## {block.content}")
                elif block.type == "equation":
                    eq = block.latex or block.content
                    parts.append(f"[EQUATION: {eq}]")
                elif block.type == "figure":
                    desc = ""
                    if block.analysis:
                        desc = block.analysis.get("description", "")
                    if not desc and block.caption:
                        desc = block.caption
                    comps = block.analysis.get("components", [])
                    comp_str = ", ".join(comps[:5]) if comps else ""
                    if comp_str:
                        parts.append(f"[FIGURE: {desc} Components: {comp_str}]")
                    else:
                        parts.append(f"[FIGURE: {desc or block.caption or 'diagram'}]")
                elif block.type == "table":
                    if block.rows:
                        rows_str = " | ".join(
                            " / ".join(cell for cell in row) for row in block.rows[:5]
                        )
                        parts.append(f"[TABLE: {rows_str}]")
                    elif block.content:
                        parts.append(f"[TABLE: {block.content}]")
                elif block.type == "caption":
                    parts.append(f"Caption: {block.content}")
                elif block.content:
                    parts.append(block.content)
        return "\n".join(parts)


# ---------------------------------------------------------------------------
# Table parser
# ---------------------------------------------------------------------------

def _parse_table_text(text: str) -> list[list[str]]:
    """Parse a pipe-separated or tab-separated table into rows."""
    rows: list[list[str]] = []
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue
        if "|" in line:
            cells = [c.strip() for c in line.split("|") if c.strip()]
            if cells:
                rows.append(cells)
        elif "\t" in line:
            cells = [c.strip() for c in line.split("\t") if c.strip()]
            if cells:
                rows.append(cells)
    return rows


# ---------------------------------------------------------------------------
# Image storage
# ---------------------------------------------------------------------------

def _save_image(image_bytes: bytes, ext: str, storage_dir: Path) -> str:
    """Save image bytes to storage directory. Returns relative path."""
    images_dir = storage_dir / "multimodal_images"
    images_dir.mkdir(parents=True, exist_ok=True)
    filename = f"{uuid4().hex}.{ext}"
    filepath = images_dir / filename
    filepath.write_bytes(image_bytes)
    return str(filepath)


# ---------------------------------------------------------------------------
# Core per-format parsers
# ---------------------------------------------------------------------------

def _parse_pdf_structured(
    content: bytes,
    storage_dir: Path,
    analyze_figures: bool = True,
    extract_math: bool = True,
) -> StructuredDocument:
    """Parse a PDF into structured pages."""
    from .pdf_extractor import extract_images_from_pdf, extract_page_renders
    from .layout_parser import detect_layout, BLOCK_TYPE_EQUATION, BLOCK_TYPE_FIGURE, BLOCK_TYPE_TABLE
    from .math_extractor import extract_math as do_extract_math
    from .figure_analyzer import analyze_figure

    # Step 1: Extract text per page using pypdf (already available)
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        page_texts: dict[int, str] = {}
        for i, page in enumerate(reader.pages):
            page_texts[i + 1] = page.extract_text() or ""
    except Exception as exc:
        logger.warning("pypdf text extraction failed: %s", exc)
        page_texts = {}

    # Step 2: Extract embedded images
    embedded_images = extract_images_from_pdf(content)
    images_by_page: dict[int, list] = {}
    for img in embedded_images:
        images_by_page.setdefault(img.page_number, []).append(img)

    # Step 3: Try page renders for layout detection (only if PyMuPDF available)
    page_renders = extract_page_renders(content, dpi=120, max_pages=30)
    renders_by_page: dict[int, bytes] = {pn: img for pn, img in page_renders}

    total_pages = max(
        len(page_texts),
        max(images_by_page.keys(), default=0),
        max(renders_by_page.keys(), default=0),
    )

    structured_pages: list[StructuredPage] = []
    equation_count = figure_count = table_count = heading_count = 0

    for page_num in range(1, total_pages + 1):
        page_text = page_texts.get(page_num, "")
        render = renders_by_page.get(page_num)
        page_images = images_by_page.get(page_num, [])

        # Layout detection
        blocks_raw = detect_layout(text=page_text, image_bytes=render)

        struct_page = StructuredPage(page_number=page_num)

        for block in blocks_raw:
            if block.block_type == BLOCK_TYPE_EQUATION and extract_math:
                math_result = do_extract_math(text=block.text)
                pb = PageBlock(
                    type="equation",
                    content=block.text,
                    latex=math_result["latex"],
                    latex_method=math_result["method"],
                    confidence=block.confidence,
                    page_number=page_num,
                    bbox=block.bbox,
                )
                equation_count += 1
                struct_page.blocks.append(pb)

            elif block.block_type == BLOCK_TYPE_TABLE:
                rows = _parse_table_text(block.text)
                pb = PageBlock(
                    type="table",
                    content=block.text,
                    rows=rows,
                    confidence=block.confidence,
                    page_number=page_num,
                    bbox=block.bbox,
                )
                table_count += 1
                struct_page.blocks.append(pb)

            elif block.block_type == BLOCK_TYPE_FIGURE and analyze_figures:
                # We'll add figure blocks from embedded images below
                pb = PageBlock(
                    type="figure",
                    content=block.text,
                    caption=block.text,
                    confidence=block.confidence,
                    page_number=page_num,
                    bbox=block.bbox,
                )
                figure_count += 1
                struct_page.blocks.append(pb)

            else:
                block_type_map = {
                    "heading":   "heading",
                    "paragraph": "paragraph",
                    "list":      "list",
                    "code":      "code",
                    "caption":   "caption",
                    "unknown":   "paragraph",
                }
                pb_type = block_type_map.get(block.block_type, "paragraph")
                if pb_type == "heading":
                    heading_count += 1
                pb = PageBlock(
                    type=pb_type,
                    content=block.text,
                    confidence=block.confidence,
                    page_number=page_num,
                    bbox=block.bbox,
                )
                struct_page.blocks.append(pb)

        # Process embedded images on this page
        for img in page_images:
            image_path = _save_image(img.image_bytes, img.extension, storage_dir)

            # Find nearby caption text
            nearby_text = ""
            for blk in struct_page.blocks:
                if blk.type in ("caption", "paragraph"):
                    nearby_text = blk.content[:300]
                    break

            analysis: dict = {}
            if analyze_figures:
                try:
                    analysis = analyze_figure(img.image_bytes, nearby_text)
                except Exception as exc:
                    logger.warning("Figure analysis failed for page %d: %s", page_num, exc)

            # Also try math extraction on image — it might be an equation
            if extract_math and img.width < img.height * 4:  # likely not an equation strip
                pb = PageBlock(
                    type="figure",
                    caption=nearby_text[:120],
                    image_path=image_path,
                    analysis=analysis,
                    confidence=0.85,
                    page_number=page_num,
                )
                figure_count += 1
                struct_page.blocks.append(pb)
            else:
                # Wide/thin image — probably an equation strip
                math_result = do_extract_math(image_bytes=img.image_bytes) if extract_math else {}
                if math_result.get("latex"):
                    pb = PageBlock(
                        type="equation",
                        latex=math_result["latex"],
                        latex_method=math_result.get("method", "vision_llm"),
                        image_path=image_path,
                        confidence=0.80,
                        page_number=page_num,
                    )
                    equation_count += 1
                    struct_page.blocks.append(pb)
                else:
                    pb = PageBlock(
                        type="figure",
                        caption=nearby_text[:120],
                        image_path=image_path,
                        analysis=analysis,
                        confidence=0.80,
                        page_number=page_num,
                    )
                    figure_count += 1
                    struct_page.blocks.append(pb)

        if struct_page.blocks:
            structured_pages.append(struct_page)

    summary = {
        "total_pages": total_pages,
        "structured_pages": len(structured_pages),
        "has_equations": equation_count > 0,
        "has_figures": figure_count > 0,
        "has_tables": table_count > 0,
        "equation_count": equation_count,
        "figure_count": figure_count,
        "table_count": table_count,
        "heading_count": heading_count,
        "parser": "pdf_structured",
    }

    return StructuredDocument(pages=structured_pages, summary=summary)


def _parse_image_structured(
    content: bytes,
    file_name: str,
    storage_dir: Path,
    analyze_figures: bool = True,
    extract_math: bool = True,
) -> StructuredDocument:
    """Parse a standalone image (PNG/JPG) into structured content."""
    from .layout_parser import detect_layout, BLOCK_TYPE_EQUATION, BLOCK_TYPE_FIGURE, BLOCK_TYPE_TABLE
    from .math_extractor import extract_math as do_extract_math
    from .figure_analyzer import analyze_figure

    ext = Path(file_name).suffix.lstrip(".").lower() or "png"
    image_path = _save_image(content, ext, storage_dir)

    # Try layout detection on the image
    blocks_raw = detect_layout(image_bytes=content)

    struct_page = StructuredPage(page_number=1)
    equation_count = figure_count = table_count = 0

    for block in blocks_raw:
        if block.block_type == BLOCK_TYPE_EQUATION and extract_math:
            math_result = do_extract_math(text=block.text, image_bytes=content)
            struct_page.blocks.append(PageBlock(
                type="equation",
                content=block.text,
                latex=math_result["latex"],
                latex_method=math_result["method"],
                confidence=block.confidence,
                page_number=1,
            ))
            equation_count += 1
        elif block.block_type == BLOCK_TYPE_TABLE:
            rows = _parse_table_text(block.text)
            struct_page.blocks.append(PageBlock(
                type="table", content=block.text, rows=rows,
                confidence=block.confidence, page_number=1,
            ))
            table_count += 1
        else:
            struct_page.blocks.append(PageBlock(
                type=block.block_type if block.block_type != "unknown" else "paragraph",
                content=block.text,
                confidence=block.confidence,
                page_number=1,
            ))

    # If we got no meaningful text blocks, treat the whole image as a figure
    text_blocks = [b for b in struct_page.blocks if b.type in ("paragraph", "heading")]
    if not text_blocks and analyze_figures:
        analysis = {}
        try:
            analysis = analyze_figure(content, "")
        except Exception as exc:
            logger.warning("Image figure analysis failed: %s", exc)
        struct_page.blocks.insert(0, PageBlock(
            type="figure",
            image_path=image_path,
            analysis=analysis,
            confidence=0.80,
            page_number=1,
        ))
        figure_count += 1

    summary = {
        "total_pages": 1,
        "structured_pages": 1,
        "has_equations": equation_count > 0,
        "has_figures": figure_count > 0,
        "has_tables": table_count > 0,
        "equation_count": equation_count,
        "figure_count": figure_count,
        "table_count": table_count,
        "heading_count": sum(1 for b in struct_page.blocks if b.type == "heading"),
        "parser": "image_structured",
    }

    return StructuredDocument(pages=[struct_page], summary=summary)


def _parse_text_document_structured(
    content: bytes,
    file_name: str,
) -> StructuredDocument:
    """Parse DOCX, PPTX, TXT, MD into structured pages using layout heuristics."""
    from .layout_parser import parse_text_layout, BLOCK_TYPE_EQUATION, BLOCK_TYPE_TABLE
    from .math_extractor import extract_math as do_extract_math

    suffix = Path(file_name).suffix.lower()

    if suffix == ".docx":
        from docx import Document as DocxDocument
        doc = DocxDocument(io.BytesIO(content))
        parts: list[str] = []
        
        # Extract images from docx relationships
        try:
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    ext = rel.target_part.content_type.split("/")[-1]
                    if ext in ["jpeg", "png", "gif", "bmp"]:
                        img_path = _save_image(rel.target_part.blob, ext, storage_dir)
                        parts.append(f"[FIGURE_PATH: {img_path}]")
        except Exception as e:
            logger.warning("Failed to extract images from DOCX: %s", e)
            
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append("| " + " | ".join(cells) + " |")
        text = "\n".join(parts)

    elif suffix in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"\n--- Slide {i} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                parts.append(para.text)
            text = "\n".join(parts)
        except Exception:
            text = content.decode("utf-8", errors="ignore")
    else:
        text = content.decode("utf-8", errors="ignore")

    blocks_raw = parse_text_layout(text)

    # Group into a single "page" — split at ~200 blocks into new pages
    PAGE_SIZE = 60
    pages_blocks: list[list] = []
    for i in range(0, max(1, len(blocks_raw)), PAGE_SIZE):
        pages_blocks.append(blocks_raw[i : i + PAGE_SIZE])

    structured_pages: list[StructuredPage] = []
    equation_count = figure_count = table_count = heading_count = 0

    for page_num, page_blocks in enumerate(pages_blocks, 1):
        struct_page = StructuredPage(page_number=page_num)
        for block in page_blocks:
            if block.block_type == BLOCK_TYPE_EQUATION:
                math_result = do_extract_math(text=block.text)
                struct_page.blocks.append(PageBlock(
                    type="equation",
                    content=block.text,
                    latex=math_result["latex"],
                    latex_method=math_result["method"],
                    confidence=block.confidence,
                    page_number=page_num,
                ))
                equation_count += 1
            elif block.block_type == BLOCK_TYPE_TABLE:
                rows = _parse_table_text(block.text)
                struct_page.blocks.append(PageBlock(
                    type="table",
                    content=block.text,
                    rows=rows,
                    confidence=block.confidence,
                    page_number=page_num,
                ))
                table_count += 1
            else:
                pb_type = block.block_type if block.block_type not in ("unknown",) else "paragraph"
                if pb_type == "heading":
                    heading_count += 1
                
                import re
                fig_match = re.search(r'\[FIGURE_PATH:\s*(.+?)\]', block.text)
                if fig_match:
                    image_path = fig_match.group(1).strip()
                    analysis = {}
                    if analyze_figures:
                        try:
                            from .figure_analyzer import analyze_figure
                            img_bytes = Path(image_path).read_bytes()
                            analysis = analyze_figure(img_bytes)
                        except Exception as e:
                            pass
                            
                    struct_page.blocks.append(PageBlock(
                        type="figure",
                        caption=f"Image {figure_count + 1}", 
                        image_path=image_path,
                        analysis=analysis,
                        confidence=block.confidence,
                        page_number=page_num,
                    ))
                    figure_count += 1
                else:
                    struct_page.blocks.append(PageBlock(
                        type=pb_type,
                        content=block.text,
                        confidence=block.confidence,
                        page_number=page_num,
                    ))
        if struct_page.blocks:
            structured_pages.append(struct_page)

    summary = {
        "total_pages": len(structured_pages),
        "structured_pages": len(structured_pages),
        "has_equations": equation_count > 0,
        "has_figures": figure_count > 0,
        "has_tables": table_count > 0,
        "equation_count": equation_count,
        "figure_count": figure_count,
        "table_count": table_count,
        "heading_count": heading_count,
        "parser": "text_structured",
    }

    return StructuredDocument(pages=structured_pages, summary=summary)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def parse_document_structure(
    file_name: str,
    content: bytes,
    storage_dir: Path,
    *,
    analyze_figures: bool = True,
    extract_math: bool = True,
) -> StructuredDocument:
    """
    Parse a document into a structured JSON representation.

    Dispatches to the appropriate parser based on file extension.

    Args:
        file_name:       Original file name (used for extension detection).
        content:         Raw file bytes.
        storage_dir:     Directory for saving extracted images.
        analyze_figures: Whether to run vision LLM figure analysis.
        extract_math:    Whether to run formula extraction on equations.

    Returns:
        StructuredDocument with per-page blocks.
    """
    suffix = Path(file_name).suffix.lower()

    logger.info("Parsing '%s' for structured content (analyze_figures=%s)", file_name, analyze_figures)

    try:
        if suffix == ".pdf":
            return _parse_pdf_structured(
                content, storage_dir,
                analyze_figures=analyze_figures,
                extract_math=extract_math,
            )
        elif suffix in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
            return _parse_image_structured(
                content, file_name, storage_dir,
                analyze_figures=analyze_figures,
                extract_math=extract_math,
            )
        elif suffix in (".docx", ".pptx", ".ppt", ".txt", ".md", ".csv"):
            return _parse_text_document_structured(content, file_name)
        else:
            # Unknown format — treat as text
            return _parse_text_document_structured(content, file_name)

    except Exception as exc:
        logger.error("Structured parsing failed for '%s': %s", file_name, exc)
        # Return minimal document so the pipeline can continue
        return StructuredDocument(
            pages=[],
            summary={
                "total_pages": 0,
                "has_equations": False,
                "has_figures": False,
                "has_tables": False,
                "equation_count": 0,
                "figure_count": 0,
                "table_count": 0,
                "heading_count": 0,
                "parser": "failed",
                "error": str(exc),
            },
        )
